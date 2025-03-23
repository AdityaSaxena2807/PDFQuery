from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from PyPDF2 import PdfReader
from pptx import Presentation
import docx
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from langchain_community.vectorstores import FAISS
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.units import inch
import io
import tempfile
import uuid
import shutil

app = Flask(__name__)
CORS(app)  # Enable CORS for all routes

# Load environment variables
load_dotenv()

# Initialize session storage
sessions = {}

# Temporary file storage directory
UPLOAD_FOLDER = os.path.join(tempfile.gettempdir(), 'PDFQuery_uploads')
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def get_session_folder(session_id):
    """Create and return a folder path for the session"""
    folder_path = os.path.join(UPLOAD_FOLDER, session_id)
    os.makedirs(folder_path, exist_ok=True)
    return folder_path

def get_pdf_text(pdf_path):
    """Extract text from PDF files"""
    text = ""
    pdf_reader = PdfReader(pdf_path)
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def get_ppt_text(ppt_path):
    """Extract text from PPT files"""
    text = ""
    presentation = Presentation(ppt_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text

def get_doc_text(doc_path):
    """Extract text from DOCX files"""
    text = ""
    doc_file = docx.Document(doc_path)
    for para in doc_file.paragraphs:
        text += para.text + " "
    return text

def get_text_chunks(text):
    """Split text into chunks for processing"""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=10000,
        chunk_overlap=1000
    )
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks, session_id):
    """Create and save a vector store for the text chunks"""
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    
    # Save to session-specific location
    index_path = os.path.join(get_session_folder(session_id), "faiss_index")
    vector_store.save_local(index_path)
    return index_path

def get_conversational_chain():
    """Create a conversational chain for answering questions"""
    prompt_template = """
    Answer the question as detailed as possible from the provided context, make sure to provide all the details, 
    if the answer is given in points make sure to give the points in different lines, 
    if the answer is not in the provided context just say, "Answer cannot be found", don't provide the wrong answer.
    Context: \n{context}\n
    Question: \n{question}\n

    Answer:  
    """
    model = ChatGoogleGenerativeAI(model="gemini-1.5-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def user_input(user_question, session_id):
    """Process user question and return answer"""
    index_path = os.path.join(get_session_folder(session_id), "faiss_index")
    
    # Check if index exists
    if not os.path.exists(index_path):
        return "Please upload and process documents first."
    
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    try:
        new_db = FAISS.load_local(index_path, embeddings, allow_dangerous_deserialization=True)
        docs = new_db.similarity_search(user_question)
        chain = get_conversational_chain()

        response = chain(
            {"input_documents": docs, "question": user_question},
            return_only_outputs=True
        )
        
        answer = response.get("output_text", "").strip()
        if not answer:
            return "Answer cannot be found"
        return answer
    except Exception as e:
        return f"Error processing question: {str(e)}"

def generate_pdf_file(qa_pairs):
    """Generate PDF from QA pairs"""
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=72)
    
    # Create styles
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=16,
        spaceAfter=30
    )
    question_style = ParagraphStyle(
        'CustomQuestion',
        parent=styles['Normal'],
        fontSize=12,
        fontName='Helvetica-Bold',
        spaceAfter=12
    )
    answer_style = ParagraphStyle(
        'CustomAnswer',
        parent=styles['Normal'],
        fontSize=12,
        fontName='Helvetica',
        spaceAfter=20,
        leftIndent=20
    )

    # Build the PDF content
    story = []
    
    # Add title
    title = Paragraph("Chat Conversation", title_style)
    story.append(title)
    
    # Add QA pairs
    for idx, (question, answer) in enumerate(qa_pairs, start=1):
        # Add question
        q_text = f"Q{idx}: {question}"
        story.append(Paragraph(q_text, question_style))
        
        # Add answer
        a_text = f"A{idx}: {answer}"
        story.append(Paragraph(a_text, answer_style))
        
        # Add space between QA pairs
        story.append(Spacer(1, 20))

    # Build the PDF
    doc.build(story)
    buffer.seek(0)
    return buffer

@app.route('/upload', methods=['POST'])
def upload_files():
    """Endpoint to upload and process files"""
    if 'files' not in request.files:
        return jsonify({"error": "No files provided"}), 400
    
    files = request.files.getlist('files')
    if not files or len(files) == 0:
        return jsonify({"error": "No files selected"}), 400
    
    # Create a new session ID if not exists in headers
    session_id = request.headers.get('X-Session-ID')
    if not session_id:
        session_id = str(uuid.uuid4())
        sessions[session_id] = {"qa_pairs": []}
    elif session_id not in sessions:
        sessions[session_id] = {"qa_pairs": []}
    
    session_folder = get_session_folder(session_id)
    
    # Save and process uploaded files
    try:
        # Clear previous files in the session folder
        for item in os.listdir(session_folder):
            file_path = os.path.join(session_folder, item)
            if os.path.isfile(file_path) and not file_path.endswith('faiss_index'):
                os.unlink(file_path)
        
        raw_text = ""
        for file in files:
            file_path = os.path.join(session_folder, file.filename)
            file.save(file_path)
            
            if file.filename.endswith('.pdf'):
                raw_text += get_pdf_text(file_path)
            elif file.filename.endswith('.pptx'):
                raw_text += get_ppt_text(file_path)
            elif file.filename.endswith('.docx'):
                raw_text += get_doc_text(file_path)
        
        # Process the combined text
        if raw_text:
            text_chunks = get_text_chunks(raw_text)
            index_path = get_vector_store(text_chunks, session_id)
            
            return jsonify({
                "message": "Files processed successfully",
                "session_id": session_id
            }), 200
        else:
            return jsonify({"error": "No text extracted from files"}), 400
            
    except Exception as e:
        return jsonify({"error": f"Error processing files: {str(e)}"}), 500

@app.route('/ask', methods=['POST'])
def ask_question():
    """Endpoint to ask questions about the uploaded documents"""
    data = request.json
    if not data or 'question' not in data:
        return jsonify({"error": "No question provided"}), 400
    
    question = data['question']
    session_id = request.headers.get('X-Session-ID')
    
    if not session_id or session_id not in sessions:
        return jsonify({"error": "Invalid session. Please upload documents first"}), 400
    
    try:
        answer = user_input(question, session_id)
        
        # Store QA pair in session
        sessions[session_id]["qa_pairs"].append((question, answer))
        
        return jsonify({
            "question": question,
            "answer": answer
        }), 200
        
    except Exception as e:
        return jsonify({"error": f"Error processing question: {str(e)}"}), 500

@app.route('/get_qa_pairs', methods=['GET'])
def get_qa_pairs():
    """Endpoint to get all QA pairs for a session"""
    session_id = request.headers.get('X-Session-ID')
    
    if not session_id or session_id not in sessions:
        return jsonify({"error": "Invalid session"}), 400
    
    return jsonify({
        "qa_pairs": sessions[session_id]["qa_pairs"]
    }), 200

@app.route('/clear_qa_pairs', methods=['POST'])
def clear_qa_pairs():
    """Endpoint to clear all QA pairs for a session"""
    session_id = request.headers.get('X-Session-ID')
    
    if not session_id or session_id not in sessions:
        return jsonify({"error": "Invalid session"}), 400
    
    sessions[session_id]["qa_pairs"] = []
    
    return jsonify({
        "message": "QA pairs cleared successfully"
    }), 200

@app.route('/generate_pdf', methods=['GET'])
def generate_pdf():
    """Endpoint to generate and download a PDF of the QA pairs"""
    session_id = request.headers.get('X-Session-ID')
    
    if not session_id or session_id not in sessions:
        return jsonify({"error": "Invalid session"}), 400
    
    qa_pairs = sessions[session_id]["qa_pairs"]
    if not qa_pairs:
        return jsonify({"error": "No QA pairs to generate PDF"}), 400
    
    try:
        pdf_buffer = generate_pdf_file(qa_pairs)
        return send_file(
            pdf_buffer,
            mimetype='application/pdf',
            as_attachment=True,
            download_name='conversations.pdf'
        )
    except Exception as e:
        return jsonify({"error": f"Error generating PDF: {str(e)}"}), 500

@app.route('/health', methods=['GET'])
def health_check():
    """Simple health check endpoint"""
    return jsonify({"status": "ok"}), 200

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)